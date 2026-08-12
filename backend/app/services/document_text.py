"""Extract plain text from common research document formats."""

from __future__ import annotations

import csv
import html
import io
import logging
import re
import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from docx import Document
from pypdf import PdfReader

logger = logging.getLogger("theresearcher")

# Hard cap so AI Checker stays local and cheap.
MAX_UPLOAD_BYTES = 12 * 1024 * 1024
MAX_EXTRACT_CHARS = 400_000
OCR_MAX_PAGES = 15
OCR_MIN_CHARS = 80

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".json",
    ".html",
    ".htm",
    ".xml",
    ".rtf",
    ".pdf",
    ".docx",
    ".pptx",
    ".odt",
    ".log",
}


class DocumentExtractError(ValueError):
    pass


def ocr_available() -> bool:
    """True when tesseract binary is on PATH (optional dependency)."""
    return bool(shutil.which("tesseract"))


def extract_text_from_upload(
    filename: str,
    data: bytes,
    content_type: str | None = None,
    *,
    enable_ocr: bool = True,
) -> dict:
    if not data:
        raise DocumentExtractError("Empty file.")
    if len(data) > MAX_UPLOAD_BYTES:
        raise DocumentExtractError(
            f"File too large (max {MAX_UPLOAD_BYTES // (1024 * 1024)} MB)."
        )

    name = (filename or "upload.bin").strip()
    ext = Path(name).suffix.lower()
    if not ext and content_type:
        ext = _ext_from_content_type(content_type)

    if ext not in SUPPORTED_EXTENSIONS:
        raise DocumentExtractError(
            f"Unsupported type '{ext or 'unknown'}'. "
            f"Use: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    ocr_used = False
    ocr_note = ""
    if ext == ".pdf":
        text, ocr_used, ocr_note = _from_pdf(data, enable_ocr=enable_ocr)
    elif ext == ".docx":
        text = _from_docx(data)
    elif ext == ".pptx":
        text = _from_pptx(data)
    elif ext == ".odt":
        text = _from_odt(data)
    elif ext in {".html", ".htm"}:
        text = _from_html(data)
    elif ext == ".rtf":
        text = _from_rtf(data)
    elif ext == ".csv":
        text = _from_csv(data)
    elif ext in {".xml"}:
        text = _from_plain(data)
    else:
        text = _from_plain(data)

    text = _normalize(text)
    if not text.strip():
        if ext == ".pdf" and enable_ocr and not ocr_available():
            raise DocumentExtractError(
                "No readable text found. This looks like a scanned PDF, but OCR is not available "
                "in this image (tesseract missing). Rebuild with OCR packages or paste text."
            )
        raise DocumentExtractError(
            "No readable text found. Scanned PDFs need OCR; try re-exporting or pasting text."
        )

    truncated = False
    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS]
        truncated = True

    return {
        "filename": name,
        "extension": ext,
        "char_count": len(text),
        "truncated": truncated,
        "text": text,
        "ocr_used": ocr_used,
        "ocr_note": ocr_note,
        "ocr_available": ocr_available(),
    }


def _ext_from_content_type(content_type: str) -> str:
    mapping = {
        "application/pdf": ".pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "application/vnd.oasis.opendocument.text": ".odt",
        "text/plain": ".txt",
        "text/markdown": ".md",
        "text/csv": ".csv",
        "text/html": ".html",
        "application/rtf": ".rtf",
        "text/rtf": ".rtf",
        "application/json": ".json",
    }
    base = (content_type or "").split(";")[0].strip().lower()
    return mapping.get(base, "")


def _from_plain(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _from_pdf(data: bytes, *, enable_ocr: bool = True) -> tuple[str, bool, str]:
    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise DocumentExtractError(f"Could not read PDF: {exc}") from exc
    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001
            continue
    text = "\n".join(parts)
    if len(text.strip()) >= OCR_MIN_CHARS:
        return text, False, ""

    if not enable_ocr:
        return text, False, "Sparse text layer; OCR disabled for this request."

    if not ocr_available():
        return text, False, "Sparse text layer; tesseract OCR not installed."

    try:
        ocr_text = _ocr_pdf_pages(data)
    except Exception as exc:  # noqa: BLE001
        logger.warning("PDF OCR failed: %s", exc)
        return text, False, f"OCR attempted but failed: {exc}"

    merged = (text + "\n" + ocr_text).strip() if text.strip() else ocr_text
    if not merged.strip():
        return "", True, "OCR ran but found no text."
    return merged, True, f"OCR used on up to {OCR_MAX_PAGES} page(s)."


def _ocr_pdf_pages(data: bytes) -> str:
    try:
        import fitz  # pymupdf
        import pytesseract
        from PIL import Image
    except ImportError as exc:
        raise DocumentExtractError(
            "OCR libraries missing (pymupdf/pytesseract/Pillow)."
        ) from exc

    doc = fitz.open(stream=data, filetype="pdf")
    parts: list[str] = []
    page_count = min(doc.page_count, OCR_MAX_PAGES)
    # Slightly lower DPI keeps OCR fast for local research packs.
    zoom = 1.6
    mat = fitz.Matrix(zoom, zoom)
    for i in range(page_count):
        page = doc.load_page(i)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        try:
            page_text = pytesseract.image_to_string(img) or ""
        except Exception as exc:  # noqa: BLE001
            logger.warning("tesseract failed on page %s: %s", i + 1, exc)
            continue
        if page_text.strip():
            parts.append(page_text)
    doc.close()
    return "\n".join(parts)


def _from_docx(data: bytes) -> str:
    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise DocumentExtractError(f"Could not read Word document: {exc}") from exc
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentExtractError("PowerPoint support is not installed.") from exc
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise DocumentExtractError(f"Could not read PowerPoint file: {exc}") from exc
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _from_odt(data: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            raw = zf.read("content.xml")
    except Exception as exc:  # noqa: BLE001
        raise DocumentExtractError(f"Could not read ODT file: {exc}") from exc
    try:
        root = ET.fromstring(raw)
    except ET.ParseError as exc:
        raise DocumentExtractError(f"Invalid ODT content: {exc}") from exc
    texts = [node.text for node in root.iter() if node.text and node.text.strip()]
    return "\n".join(texts)


def _from_html(data: bytes) -> str:
    raw = _from_plain(data)
    raw = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", raw)
    raw = re.sub(r"(?s)<[^>]+>", " ", raw)
    return html.unescape(raw)


def _from_rtf(data: bytes) -> str:
    raw = _from_plain(data)
    # Lightweight RTF stripper for common plain research exports.
    raw = re.sub(r"\\'[0-9a-fA-F]{2}", " ", raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d*[ ]?", " ", raw)
    raw = raw.replace("{", " ").replace("}", " ")
    return raw


def _from_csv(data: bytes) -> str:
    text = _from_plain(data)
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
    return "\n".join(r for r in rows if r)


def _normalize(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
